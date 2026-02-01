from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def set_text_style(shape, color, size, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color
            run.font.size = size
            run.font.bold = bold
            run.font.name = 'Arial'

def add_sidebar_accent(slide):
    # Adds a subtle vertical bar on the left for a high-tech vibe
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(37, 99, 235) # Blue-600
    shape.line.transparent = True

def create_premium_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # Widescreen
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(15, 23, 42)    # Dark Navy (Dashboard Bg)
    title_color = RGBColor(255, 255, 255) # White
    accent_color = RGBColor(96, 165, 250) # Blue-400
    text_color = RGBColor(203, 213, 225) # Slate-300

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    apply_background(slide, bg_color)
    add_sidebar_accent(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GLOSA-BHARAT"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(80)
    run.font.bold = True
    run.font.color.rgb = title_color

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(4), Inches(10), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Intelligent Traffic Flow Optimization for Indian Smart Cities\nAI for Atmanirbhar Bharat Initiative"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.color.rgb = accent_color

    # Slide 2: Problem
    def add_content_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, bg_color)
        add_sidebar_accent(slide)
        
        # Header bar
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0.5), Inches(11.5), Inches(0.8))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
        header_bar.line.transparent = True
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.6), Inches(10), Inches(0.6))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        run = p.runs[0]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = accent_color
        
        # Content
        c_box = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(11), Inches(4))
        tf = c_box.text_frame
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.space_after = Pt(15)
            run = p.runs[0]
            run.font.size = Pt(24)
            run.font.color.rgb = text_color

    slides_content = [
        ("The Challenge: Urban Gridlock", [
            "• Congestion costs India ₹2.4 Lakh Crore annually.",
            "• Massive fuel wastage and CO2 emissions from idling.",
            "• Conventional traffic signals are dynamic-blind (fixed timers).",
            "• Need for an indigenous, AI-driven traffic sync protocol."
        ]),
        ("Our Solution: GLOSA-BHARAT", [
            "• Green Light Optimal Speed Advisory (GLOSA).",
            "• AI analyzes real-time junction density & signal state.",
            "• Advises drivers on the 'Perfect Speed' to catch the green light.",
            "• Eliminates 'Stop-and-Go' behavior for smoother mobility."
        ]),
        ("System Architecture", [
            "• PERCEPTION: YOLOv8 vehicle detection & Queue length analysis.",
            "• BRAIN: Node.js coordination engine for signal-vehicle telemetry.",
            "• INTERFACE: Real-time React dashboard with Leaflet GIS mapping.",
            "• EDGE: V2I synchronization with sub-second latency."
        ]),
        ("The Technology Stack", [
            "• AI/Vision: Python, OpenCV, Darknet YOLO.",
            "• Core Backend: Node.js, Express, Axios.",
            "• Modern Frontend: React.js, Tailwind CSS, Framer Motion.",
            "• Data Flow: Real-time WebSockets for telemetry updates."
        ]),
        ("Innovation for Indian Roads", [
            "• MIXED TRAFFIC: Trained for bikes, autos, and large fleets.",
            "• SCALABLE: Designed for quick rollout in Tier-2/3 cities.",
            "• ATMANIRBHAR: 100% indigenous software stack.",
            "• NO PROPRIETARY LOCK-IN: Built on open-source standards."
        ]),
        ("Impact & Strategic Value", [
            "• ECONOMY: 15-20% reduction in delivery fleet fuel costs.",
            "• ENVIRONMENT: Drastic drop in 'Cold Start' idling emissions.",
            "• SOCIAL: Reduced driver stress and predictable arrival times.",
            "• VISION: A smarter, self-reliant mobility infrastructure."
        ])
    ]

    for title, bullets in slides_content:
        add_content_slide(title, bullets)

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, bg_color)
    add_sidebar_accent(slide)
    t_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SMARTER CITIES. SELF-RELIANT INDIA."
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = title_color
    
    p = tf.add_paragraph()
    p.text = "THANK YOU | Q&A"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = accent_color

    prs.save('GLOSA_Bharat_Premium_Presentation.pptx')
    print("Premium PPT generated successfully!")

if __name__ == "__main__":
    create_premium_presentation()
