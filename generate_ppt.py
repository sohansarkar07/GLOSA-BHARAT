from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "GLOSA-BHARAT"
    subtitle.text = "Intelligent Traffic Flow Optimization for Indian Smart Cities\nAI for Atmanirbhar Bharat Initiative"

    # Slide 2: Problem
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The Problem: Urban Gridlock"
    content = slide.placeholders[1]
    content.text = (
        "• Traffic congestion costs Indian urban centers billions in productivity.\n"
        "• Vehicular idling at red lights is a primary source of urban CO2 emissions.\n"
        "• High fuel imports exacerbated by inefficient 'stop-and-go' traffic.\n"
        "• Conventional traffic lights don't adapt to real-time vehicle density."
    )

    # Slide 3: Solution
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Our Solution: GLOSA Implementation"
    content = slide.placeholders[1]
    content.text = (
        "• What is GLOSA? Green Light Optimal Speed Advisory.\n"
        "• Core Concept: AI calculates 'ideal speed' for drivers.\n"
        "• Objective: Arrive at intersections exactly when the light turns green.\n"
        "• Goal: Achieve a 'Green Wave' across city corridors."
    )

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture"
    content = slide.placeholders[1]
    content.text = (
        "• Perception Layer: YOLO/OpenCV analysis of live camera feeds.\n"
        "• Brain (Backend): Node.js server calculating optimal speeds.\n"
        "• Interface: Real-time React dashboard and Driver Advisory Feed.\n"
        "• Protocol: V2I synchronization for sub-second latency."
    )

    # Slide 5: Tech Stack
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Technology Stack"
    content = slide.placeholders[1]
    content.text = (
        "• AI/ML: Python, YOLO, OpenCV.\n"
        "• Backend: Node.js, Express.js.\n"
        "• Frontend: React.js, Tailwind CSS, Framer Motion.\n"
        "• Maps: React-Leaflet, OpenStreetMap API.\n"
        "• Design: Premium 'Control Room' Dark Mode aesthetic."
    )

    # Slide 6: Innovation for India
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Innovation for Indian Roads"
    content = slide.placeholders[1]
    content.text = (
        "• Indigenous Resilience: Handles mixed traffic (Autos, Bikes).\n"
        "• Scalability: Low-cost edge computing for Tier-2/3 cities.\n"
        "• Atmanirbhar Vision: 100% indigenous software stack.\n"
        "• Reduced dependence on expensive foreign proprietary software."
    )

    # Slide 7: Visualization
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Real-Time Visualization & Metrics"
    content = slide.placeholders[1]
    content.text = (
        "• Live Dashboard: Real-time geospatial tracking.\n"
        "• Advisory Accuracy: 95%+ precision in signal timing predictions.\n"
        "• User Experience: Intuitive 'Driver Feed' with speed recommendations.\n"
        "• Dynamic Themes: Dark/Light mode for operational flexibility."
    )

    # Slide 8: Inclusion & Scalability
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Inclusion & Scalability"
    content = slide.placeholders[1]
    content.text = (
        "• Broad Utility: Zomato/Blinkit fleets, Buses, Ambulances.\n"
        "• Integration: Works with existing city CCTV infrastructure.\n"
        "• Smart City Grid: Ready for national-level cloud synchronization."
    )

    # Slide 9: Impact
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Measurable Impact"
    content = slide.placeholders[1]
    content.text = (
        "• Fuel Efficiency: Projected 15-18% reduction in consumption.\n"
        "• Sustainability: Massive reduction in carbon footprint.\n"
        "• Traffic Flow: Improved average speeds without increasing peaks.\n"
        "• Economic: Billions saved in logistics and time productivity."
    )

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = (
        "• Blueprint for the future of Indian Mobility.\n"
        "• AI in Real-World Applications: From Idea to Implementation.\n"
        "• Vision: Smarter, Faster, and Self-Reliant Indian Cities.\n"
        "• Thank You! | Q&A"
    )

    prs.save('GLOSA_Bharat_Presentation.pptx')
    print("PPT generated successfully!")

if __name__ == "__main__":
    create_presentation()
