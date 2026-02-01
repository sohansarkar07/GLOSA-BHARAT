from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def create_notebook_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # NotebookLM Colors
    bg_color = RGBColor(248, 250, 252)    # Very light sky/grey
    card_color = RGBColor(255, 255, 255)  # White
    accent_teal = RGBColor(13, 148, 136)  # Teal-600
    text_main = RGBColor(15, 23, 42)      # Slate-900
    text_sub = RGBColor(71, 85, 105)       # Slate-600

    def add_notebook_slide(title_text, subtitle_text, content_points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, bg_color)
        
        # Left Panel (Source Indicator)
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), Inches(7.5))
        panel.fill.solid()
        panel.fill.fore_color.rgb = accent_teal
        panel.line.transparent = True

        # Header Area
        t_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = text_main
        
        s_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(0.5))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        run = p.runs[0]
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = text_sub

        # Main Card Content
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.5), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.color.rgb = RGBColor(226, 232, 240)
        card.line.width = Pt(1)

        # Content in Card
        c_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.5), Inches(4))
        tf = c_box.text_frame
        tf.word_wrap = True
        for idx, point in enumerate(content_points):
            p = tf.add_paragraph()
            p.text = f"[{idx+1}] Source Intelligence: {point}"
            p.space_after = Pt(20)
            run = p.runs[0]
            run.font.size = Pt(20)
            run.font.color.rgb = text_main
            run.font.name = 'Arial'

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, card_color)
    
    # Large Teal Accent Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(8), Inches(8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_teal
    circle.line.transparent = True
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GLOSA-BHARAT"
    run = p.runs[0]
    run.font.size = Pt(64)
    run.font.bold = True
    run.font.color.rgb = text_main
    
    p2 = tf.add_paragraph()
    p2.text = "Analyzing Smart Mobility for Atmanirbhar Bharat"
    run2 = p2.runs[0]
    run2.font.size = Pt(24)
    run2.font.color.rgb = text_sub

    # Slides content
    slides_data = [
        ("The Congestion Thesis", "Analyzing Economic and Environmental Friction", [
            "Urban gridlock costs India ₹2.4 Lakh Crore annually in lost productivity.",
            "Vehicle idling at junctions contributes to 25% of localized urban emissions.",
            "Conventional timed signals act as bottlenecks rather than flow-facilitators.",
            "Data suggests 18% of urban fuel is wasted in stop-and-go behavior."
        ]),
        ("Strategic Implementation", "Synthesizing AI and Infrastructure", [
            "GLOSA acts as a proactive speed advisory layer for connected vehicles.",
            "AI analyzes queue length and signal state to predict intersection availability.",
            "Drivers receive a calculated 'Perfect Speed' to minimize stopping.",
            "Increases road capacity without requiring new physical lanes."
        ]),
        ("Technical Insights", "System Architecture & Data Telemetry", [
            "Computer Vision: YOLOv8 model trained on heterogeneous Indian traffic data.",
            "Orchestration: Node.js backend handles sub-second latency telemetry maps.",
            "Interaction: React-based dashboard visualizing real-time flow and analytics.",
            "Network: V2X (Vehicle-to-Everything) synchronization framework."
        ]),
        ("The 'Atmanirbhar' Edge", "Evaluating Localized Innovation", [
            "100% Indigenous software stack tailored for Indian road behavior.",
            "Optimized for low-bandwidth environments through efficient JSON messaging.",
            "Compatible with existing Smart City CCTV and sensor infrastructure.",
            "Reduces capital expenditure by over 40% compared to Western solutions."
        ])
    ]

    for title, sub, points in slides_data:
        add_notebook_slide(title, sub, points)

    # Conclusion Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, bg_color)
    
    fin_box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1.5))
    tf = fin_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Intelligent Mobility synthesized for India."
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = text_main
    
    p2 = tf.add_paragraph()
    p2.text = "Q&A | Research Summary Complete"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(18)
    run2.font.color.rgb = accent_teal

    prs.save('GLOSA_NotebookLM_Presentation.pptx')
    print("NotebookLM style PPT generated successfully!")

if __name__ == "__main__":
    create_notebook_ppt()
