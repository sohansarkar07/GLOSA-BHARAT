from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_sidebar_accent(slide, color_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.1), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_rgb
    shape.line.transparent = True

def create_hacker_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Hackathon Palette
    bg_dark = RGBColor(15, 23, 42)
    accent_blue = RGBColor(59, 130, 246)
    accent_emerald = RGBColor(16, 185, 129)
    text_white = RGBColor(248, 250, 252)

    # SLIDE 1: THE HOOK
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, bg_dark)
    add_sidebar_accent(slide, accent_blue)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GLOSA-BHARAT"
    run = p.runs[0]
    run.font.size = Pt(88)
    run.font.bold = True
    run.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "Solving Urban Friction with Indigenous AI"
    run2 = p2.runs[0]
    run2.font.size = Pt(32)
    run2.font.color.rgb = accent_blue

    # SLIDE 2: THE PROBLEM (BIG NUMBERS)
    def add_hero_slide(title, subtitle, bullets, accent=accent_blue):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, bg_dark)
        add_sidebar_accent(slide, accent)
        
        # Header
        t_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        p = t_box.text_frame.paragraphs[0]
        p.text = title.upper()
        run = p.runs[0]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = text_white
        
        s_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(0.5))
        p = s_box.text_frame.paragraphs[0]
        p.text = subtitle
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = accent
        
        # Content
        c_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
        tf = c_box.text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.space_after = Pt(20)
            run = p.runs[0]
            run.font.size = Pt(28)
            run.font.color.rgb = text_white

    slides = [
        ("The Challenge: Urban Friction", "The Trillion Rupee Problem", [
            "• Congestion costs India ₹2.4 Lakh Crore / Year.",
            "• Massive fuel wastage due to 'Stop-and-Go' behavior.",
            "• Signals are 'Timed-Blind' to real-time traffic flow.",
            "• High localized emissions at signal hotspots."
        ]),
        ("System Architecture", "The Engineering Ecosystem", [
            "• PERCEPTION: Edge AI (YOLOv8) for real-time ID.",
            "• ORCHESTRATION: High-concurrency Node.js Engine.",
            "• ADVISORY: Mathematical GLOSA Optimization.",
            "• INTERFACE: React Digital Twin Map Dashboard."
        ]),
        ("Why We Win (Innovation)", "The Sovereign Advantage", [
            "• INDIGENOUS: Trained for Autos, Bikes & Mixed traffic.",
            "• ACCESSIBLE: Uses existing infra (No LIDAR needed).",
            "• REAL-TIME: WebSockets ensure <100ms latency.",
            "• SCALABLE: Built on modular Microservices Stack."
        ]),
        ("The 'Aatmanirbhar' Impact", "Self-Reliant Smart Cities", [
            "• Economic: 15-20% fuel savings for logistics fleets.",
            "• Strategic: 100% Homemade, Secure Traffic Stack.",
            "• Environmental: Drastic reduction in carbon footprint.",
            "• Vision: Ready for Tier-2 & Tier-3 Smart City rollout."
        ])
    ]

    for t, s, b in slides:
        add_hero_slide(t, s, b)

    # FINAL SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, bg_dark)
    add_sidebar_accent(slide, accent_emerald)
    
    end_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Intelligent Roads. Self-Reliant India."
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "GLOSA-BHARAT | HACKATHON 2026"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(24)
    run2.font.color.rgb = accent_emerald

    prs.save('GLOSA_HackerEdition_Presentation.pptx')
    print("Hacker Edition PPT generated successfully!")

if __name__ == "__main__":
    create_hacker_ppt()
